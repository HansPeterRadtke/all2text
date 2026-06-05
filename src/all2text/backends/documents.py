from __future__ import annotations

import re
import zipfile
from pathlib import Path
from typing import Any
from xml.etree import ElementTree as ET

from all2text.backends.binary import binary_summary_text
from all2text.config import config_for_context
from all2text.models import Classification, ConversionContext, ConversionResult

WORD_NS = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
ODF_NS = {
    "office": "urn:oasis:names:tc:opendocument:xmlns:office:1.0",
    "text": "urn:oasis:names:tc:opendocument:xmlns:text:1.0",
    "table": "urn:oasis:names:tc:opendocument:xmlns:table:1.0",
}
_CROSS_SHEET_REF_RE = re.compile(r"(?:'[^']+'|[A-Za-z0-9_]+)![$A-Z0-9:]+")


class DocumentBackend:
    name = "document_native_backend"

    CATEGORIES = {"document", "spreadsheet", "presentation"}

    def can_handle(self, classification: Classification, entry_type: str) -> bool:
        return entry_type == "file" and classification.rough_category in self.CATEGORIES

    def convert(
        self,
        path: Path,
        rel_path: Path,
        classification: Classification,
        metadata: dict[str, object],
        ctx: ConversionContext,
    ) -> ConversionResult:
        fmt = classification.concrete_format.upper()
        if fmt in {"DOCX", "XLSX", "PPTX", "PDF", "XLS"} and (
            not ctx.options.allow_optional_python or not ctx.options.auto_detect_python
        ):
            error = (
                f"optional_python_disabled_by_profile:{ctx.options.profile}"
                if not ctx.options.allow_optional_python
                else "optional_python_disabled_by_run.auto_detect_python=false"
            )
            return document_fallback(
                path,
                classification,
                ctx,
                dependency_error=error,
            )
        if fmt == "DOCX":
            return convert_docx(path, rel_path, classification, ctx)
        if fmt == "XLSX":
            return convert_xlsx(path, rel_path, classification, ctx)
        if fmt == "PPTX":
            return convert_pptx(path, rel_path, classification, ctx)
        if fmt == "PDF":
            return convert_pdf(path, rel_path, classification, ctx)
        if fmt in {"ODT", "ODS", "ODP"}:
            return convert_odf(path, rel_path, classification, ctx)
        if fmt == "XLS":
            return convert_xls(path, rel_path, classification, ctx)
        return document_fallback(path, classification, ctx)


DocumentPlaceholderBackend = DocumentBackend


def convert_docx(
    path: Path,
    rel_path: Path,
    classification: Classification,
    ctx: ConversionContext,
) -> ConversionResult:
    try:
        from docx import Document
        from docx.document import Document as _Document
        from docx.oxml.ns import qn
        from docx.oxml.table import CT_Tbl
        from docx.oxml.text.paragraph import CT_P
        from docx.table import _Cell, Table
        from docx.text.paragraph import Paragraph
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"python-docx unavailable:{exc}")

    try:
        document = Document(str(path))
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"python-docx open failed:{exc}")

    raw_paragraph_count = raw_docx_paragraph_count(path)
    visible_paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    ordered_blocks: list[str] = []
    tables: list[list[list[str]]] = []
    image_count = 0

    def iter_block_items(parent: Any) -> Any:
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise TypeError(parent)
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def paragraph_hyperlinks(paragraph: Any) -> list[str]:
        links: list[str] = []
        for link in paragraph._element.xpath(".//w:hyperlink"):
            rid = link.get(qn("r:id"))
            if rid and rid in paragraph.part.rels:
                links.append(paragraph.part.rels[rid].target_ref)
            anchor = link.get(qn("w:anchor"))
            if anchor:
                links.append(f"#{anchor}")
        return links

    def paragraph_image_count(paragraph: Any) -> int:
        count = 0
        for blip in paragraph._element.xpath(".//a:blip"):
            rid = blip.get(qn("r:embed"))
            if rid and document.part.related_parts.get(rid) is not None:
                count += 1
        return count

    for block_index, block in enumerate(iter_block_items(document), start=1):
        if isinstance(block, Paragraph):
            text = block.text.strip()
            style = getattr(getattr(block, "style", None), "name", None)
            hyperlinks = paragraph_hyperlinks(block)
            images = paragraph_image_count(block)
            image_count += images
            ordered_blocks.append(
                f"Block {block_index}: paragraph; style={style!r}; text={text!r}; "
                f"hyperlinks={hyperlinks or 'none'}; embedded_images={images}"
            )
        elif isinstance(block, Table):
            table_rows: list[list[str]] = []
            for row in block.rows:
                table_rows.append(
                    [
                        " / ".join(paragraph.text.strip() for paragraph in cell.paragraphs if paragraph.text.strip())
                        for cell in row.cells
                    ]
                )
            tables.append(table_rows)
            table_index = len(tables)
            ordered_blocks.append(
                f"Block {block_index}: table_{table_index}; rows={len(table_rows)}; "
                f"columns={max((len(row) for row in table_rows), default=0)}"
            )
            for row_index, row in enumerate(table_rows, start=1):
                ordered_blocks.append(f"Block {block_index}: table_{table_index} row {row_index}: " + " | ".join(row))

    lines = [
        f"File: {rel_path.name}",
        "Description: DOCX extraction using python-docx plus raw WordprocessingML counts.",
        "",
        "Document metadata:",
        *docx_core_property_lines(document),
        f"- visible_paragraph_count: {len(visible_paragraphs)}",
        f"- raw_wordprocessingml_paragraph_count: {raw_paragraph_count}",
        f"- table_count: {len(tables)}",
        f"- embedded_image_count: {image_count}",
        f"- section_count: {len(document.sections)}",
    ]
    for index, section in enumerate(document.sections, start=1):
        lines.extend(
            [
                f"Section {index}:",
                f"- page_width: {section.page_width}",
                f"- page_height: {section.page_height}",
                f"- left_margin: {section.left_margin}",
                f"- right_margin: {section.right_margin}",
                f"- header_text: {paragraph_collection_text(section.header.paragraphs)}",
                f"- footer_text: {paragraph_collection_text(section.footer.paragraphs)}",
            ]
        )
    lines.extend(["", "Visible paragraphs in document order:"])
    lines.extend([f"- paragraph {index}: {text!r}" for index, text in enumerate(visible_paragraphs, start=1)] or ["- none"])
    lines.extend(["", "Ordered block extraction:"])
    lines.extend(f"- {block}" for block in ordered_blocks)
    return ConversionResult(
        text="\n".join(lines).rstrip() + "\n",
        converter_used="docx_native_backend",
        extraction_methods_used=["python_docx_ordered_blocks", "wordprocessingml_raw_count"],
        metadata={
            "format": "DOCX",
            "visible_paragraph_count": len(visible_paragraphs),
            "raw_wordprocessingml_paragraph_count": raw_paragraph_count,
            "table_count": len(tables),
            "embedded_image_count": image_count,
        },
        limitations=["Embedded images are counted and located; image OCR/VLM analysis is routed by image providers, not faked here."]
        if image_count
        else [],
    )


def convert_xlsx(
    path: Path,
    rel_path: Path,
    classification: Classification,
    ctx: ConversionContext,
) -> ConversionResult:
    params = document_module_params(ctx, classification)
    include_hidden_sheets = bool_param(params.get("include_hidden_sheets"), True)
    max_cells_per_sheet = positive_int(params.get("max_cells_per_sheet"))
    try:
        from openpyxl import load_workbook
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"openpyxl unavailable:{exc}")
    try:
        workbook = load_workbook(path, data_only=False)
        workbook_values = load_workbook(path, data_only=True)
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"openpyxl open failed:{exc}")

    lines: list[str] = [
        f"File: {rel_path.name}",
        "Description: XLSX extraction using openpyxl. Worksheets, non-empty cells, formulas, cached values, tables, and chart metadata are listed.",
        "",
        "Workbook:",
        f"- worksheet_count: {len(workbook.worksheets)}",
        "- sheet_order: " + ", ".join(ws.title for ws in workbook.worksheets),
        *workbook_properties_lines(workbook),
    ]
    lines.extend(defined_names_lines(workbook))
    total_cells = 0
    formula_count = 0
    chart_count = 0
    hidden_sheets: list[str] = []
    skipped_hidden_sheets: list[str] = []
    truncated_sheets: list[str] = []
    for sheet_index, sheet in enumerate(workbook.worksheets, start=1):
        value_sheet = workbook_values[sheet.title]
        if sheet.sheet_state != "visible":
            hidden_sheets.append(sheet.title)
            if not include_hidden_sheets:
                skipped_hidden_sheets.append(sheet.title)
                lines.extend(
                    [
                        "",
                        f"Worksheet {sheet_index}: {sheet.title}",
                        "- visibility: hidden",
                        "- extraction: skipped_by_module_config",
                    ]
                )
                continue
        sheet_cells, truncated = non_empty_cells(sheet, max_cells=max_cells_per_sheet)
        if truncated:
            truncated_sheets.append(sheet.title)
        total_cells += len(sheet_cells)
        formula_count += sum(1 for cell in sheet_cells if is_formula(cell.value))
        chart_count += len(getattr(sheet, "_charts", []) or [])
        lines.extend(
            xlsx_sheet_lines(
                sheet_index,
                sheet,
                value_sheet,
                workbook_values,
                sheet_cells,
                truncated=truncated,
            )
        )
    limitations: list[str] = []
    warnings: list[str] = []
    if skipped_hidden_sheets:
        limitations.append("Hidden worksheets were skipped because spreadsheet.include_hidden_sheets=false.")
        warnings.append("xlsx_hidden_sheets_skipped:" + ",".join(skipped_hidden_sheets))
    if truncated_sheets:
        limitations.append("Worksheet cell extraction was limited by spreadsheet.max_cells_per_sheet.")
        warnings.append("xlsx_cells_truncated:" + ",".join(truncated_sheets))
    return ConversionResult(
        text="\n".join(lines).rstrip() + "\n",
        converter_used="xlsx_native_backend",
        extraction_methods_used=[
            "openpyxl_workbook_metadata",
            "openpyxl_all_sheets",
            "openpyxl_all_non_empty_cells",
            "openpyxl_formula_view",
            "openpyxl_cached_value_view",
            "openpyxl_chart_metadata",
        ],
        metadata={
            "format": "XLSX",
            "sheet_names": [ws.title for ws in workbook.worksheets],
            "hidden_sheets": hidden_sheets,
            "skipped_hidden_sheets": skipped_hidden_sheets,
            "non_empty_cell_count": total_cells,
            "formula_count": formula_count,
            "chart_count": chart_count,
            "max_cells_per_sheet": max_cells_per_sheet,
            "truncated_sheets": truncated_sheets,
        },
        warnings=warnings,
        limitations=limitations,
    )


def convert_pptx(
    path: Path,
    rel_path: Path,
    classification: Classification,
    ctx: ConversionContext,
) -> ConversionResult:
    try:
        from pptx import Presentation
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"python-pptx unavailable:{exc}")
    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"python-pptx open failed:{exc}")

    lines: list[str] = [
        f"File: {rel_path.name}",
        "Description: PPTX extraction using python-pptx with slide text, notes, shape metadata, and embedded image counts.",
        "",
        f"slide_count: {len(presentation.slides)}",
        f"slide_width: {presentation.slide_width}",
        f"slide_height: {presentation.slide_height}",
    ]
    image_count = 0
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.extend(["", f"Slide {slide_index}:", f"- shape_count: {len(slide.shapes)}"])
        for shape_index, shape in enumerate(slide.shapes, start=1):
            shape_type = getattr(shape, "shape_type", None)
            lines.append(
                f"- shape {shape_index}: type={shape_type}; name={getattr(shape, 'name', None)!r}; "
                f"left={getattr(shape, 'left', None)}; top={getattr(shape, 'top', None)}; "
                f"width={getattr(shape, 'width', None)}; height={getattr(shape, 'height', None)}"
            )
            if getattr(shape, "has_text_frame", False):
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs, start=1):
                    text = paragraph.text.strip()
                    if text:
                        lines.append(f"  - paragraph {paragraph_index}: {text!r}")
            if shape_type == MSO_SHAPE_TYPE.PICTURE and getattr(shape, "image", None) is not None:
                image_count += 1
                lines.append(f"  - embedded_image_{image_count}: content_type={shape.image.content_type!r}")
        notes = slide_notes(slide)
        lines.append(f"- notes: {notes if notes else 'none'}")
    return ConversionResult(
        text="\n".join(lines).rstrip() + "\n",
        converter_used="pptx_native_backend",
        extraction_methods_used=["python_pptx_shapes", "python_pptx_notes"],
        metadata={"format": "PPTX", "slide_count": len(presentation.slides), "embedded_image_count": image_count},
        limitations=["Embedded images are counted; image OCR/VLM analysis is routed by image providers, not faked here."]
        if image_count
        else [],
    )


def convert_pdf(
    path: Path,
    rel_path: Path,
    classification: Classification,
    ctx: ConversionContext,
) -> ConversionResult:
    params = document_module_params(ctx, classification)
    max_pdf_pages = positive_int(params.get("max_pdf_pages"))
    try:
        from pypdf import PdfReader
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"pypdf unavailable:{exc}")
    try:
        reader = PdfReader(str(path), strict=False)
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"pypdf open failed:{exc}")

    lines: list[str] = [
        f"File: {rel_path.name}",
        "Description: PDF extraction using pypdf native text. OCR/rendering is only attempted by configured providers.",
        "",
        f"PDF metadata: {pdf_metadata_text(reader)}",
        f"Page count: {len(reader.pages)}",
    ]
    total_chars = 0
    image_count = 0
    methods = ["pypdf_text_extraction"]
    page_count = len(reader.pages)
    extract_count = min(page_count, max_pdf_pages) if max_pdf_pages else page_count
    skipped_pages = max(0, page_count - extract_count)
    for page_number in range(1, extract_count + 1):
        page = reader.pages[page_number - 1]
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            text = ""
            lines.append(f"- page_{page_number}_extract_text_error: {exc}")
        page_images = len(getattr(page, "images", []) or [])
        image_count += page_images
        total_chars += len(text)
        lines.extend(["", f"Page {page_number}:", f"- native_text_character_count: {len(text)}", f"- pypdf_image_count: {page_images}"])
        if text.strip():
            lines.append("- native_text:")
            lines.extend(f"  {line}" for line in text.splitlines())
        else:
            lines.append("- native_text: <none>")
            lines.append("- ocr_status: skipped_no_configured_pdf_page_renderer_or_ocr_provider")
    limitations = []
    if total_chars == 0:
        limitations.append("No native PDF text was found; scanned-page OCR requires configured OCR/rendering providers.")
    if image_count:
        limitations.append("Embedded PDF images are counted; deep image analysis is not claimed unless image providers run.")
    if skipped_pages:
        limitations.append("PDF page extraction was limited by document.max_pdf_pages.")
        lines.extend(["", f"Skipped pages due to max_pdf_pages: {skipped_pages}"])
    return ConversionResult(
        text="\n".join(lines).rstrip() + "\n",
        converter_used="pdf_native_backend",
        extraction_methods_used=methods,
        metadata={
            "format": "PDF",
            "page_count": page_count,
            "extracted_page_count": extract_count,
            "skipped_page_count": skipped_pages,
            "native_text_characters": total_chars,
            "image_count": image_count,
            "max_pdf_pages": max_pdf_pages,
        },
        limitations=limitations,
    )


def convert_odf(
    path: Path,
    rel_path: Path,
    classification: Classification,
    ctx: ConversionContext,
) -> ConversionResult:
    params = document_module_params(ctx, classification)
    max_text_blocks = positive_int(params.get("max_text_blocks"))
    try:
        with zipfile.ZipFile(path) as archive:
            content = archive.read("content.xml")
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"OpenDocument content.xml unavailable:{exc}")
    try:
        root = ET.fromstring(content)
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"OpenDocument XML parse failed:{exc}")
    fmt = classification.concrete_format.upper()
    text_blocks, truncated = odf_text_blocks(root, max_blocks=max_text_blocks)
    lines = [
        f"File: {rel_path.name}",
        f"Description: {fmt} OpenDocument content.xml text extraction.",
        "",
        f"text_block_count: {len(text_blocks)}",
        "",
        "Text blocks:",
    ]
    lines.extend([f"- {block}" for block in text_blocks] or ["- none"])
    return ConversionResult(
        text="\n".join(lines).rstrip() + "\n",
        converter_used="opendocument_native_backend",
        extraction_methods_used=["zip_content_xml_parse", "opendocument_text_nodes"],
        metadata={
            "format": fmt,
            "text_block_count": len(text_blocks),
            "max_text_blocks": max_text_blocks,
            "truncated": truncated,
        },
        limitations=[
            "OpenDocument extraction is text-node oriented; layout, styles, formulas, and embedded media are not deeply extracted.",
            *(
                ["OpenDocument text blocks were limited by document.max_text_blocks."]
                if truncated
                else []
            ),
        ],
    )


def convert_xls(
    path: Path,
    rel_path: Path,
    classification: Classification,
    ctx: ConversionContext,
) -> ConversionResult:
    try:
        import xlrd  # type: ignore[import-not-found]
    except Exception as exc:
        return document_fallback(
            path,
            classification,
            ctx,
            dependency_error=f"xlrd unavailable for legacy XLS extraction:{exc}",
        )
    try:
        workbook = xlrd.open_workbook(str(path), on_demand=True)
    except Exception as exc:
        return document_fallback(path, classification, ctx, dependency_error=f"xlrd open failed:{exc}")
    lines = [f"File: {rel_path.name}", "Description: Legacy XLS extraction using xlrd.", "", f"sheet_count: {workbook.nsheets}"]
    cell_count = 0
    for sheet in workbook.sheets():
        lines.extend(["", f"Worksheet: {sheet.name}", f"- rows: {sheet.nrows}", f"- columns: {sheet.ncols}", "- cells:"])
        for row_index in range(sheet.nrows):
            for col_index in range(sheet.ncols):
                value = sheet.cell_value(row_index, col_index)
                if value not in ("", None):
                    cell_count += 1
                    lines.append(f"  - row={row_index + 1}; column={col_index + 1}; value={value!r}")
    return ConversionResult(
        text="\n".join(lines).rstrip() + "\n",
        converter_used="xls_native_backend",
        extraction_methods_used=["xlrd_workbook_read"],
        metadata={"format": "XLS", "sheet_count": workbook.nsheets, "non_empty_cell_count": cell_count},
        limitations=["Legacy XLS formulas/styles are not deeply extracted by the core xlrd path."],
    )


def document_fallback(
    path: Path,
    classification: Classification,
    ctx: ConversionContext,
    *,
    dependency_error: str | None = None,
) -> ConversionResult:
    limitation = (
        "Native extraction for this document format was not available or failed. "
        "Core all2text recorded safe metadata and did not claim semantic extraction."
    )
    doc_meta = document_light_metadata(path, classification)
    extra = [f"- limitation: {limitation}"]
    warnings: list[str] = []
    if dependency_error:
        extra.append(f"- native_extraction_status: {dependency_error}")
        warnings.append(dependency_error)
    if doc_meta:
        extra.append("- document_metadata: " + repr(doc_meta))
    return ConversionResult(
        text=binary_summary_text(path, classification, ctx, heading="Document safe summary", extra_lines=extra),
        converter_used="document_placeholder_backend",
        extraction_methods_used=["document_placeholder_summary"],
        warnings=warnings,
        metadata={"document": doc_meta, "native_extraction_status": dependency_error or "unsupported_format"},
        limitations=[limitation],
    )


def document_light_metadata(path: Path, classification: Classification) -> dict[str, object]:
    fmt = classification.concrete_format.upper()
    result: dict[str, object] = {"format": classification.concrete_format}
    if fmt == "PDF":
        try:
            raw = path.read_bytes()
            result["pdf_header"] = raw[:16].decode("ascii", errors="replace")
            result["rough_page_marker_count"] = len(re.findall(rb"/Type\s*/Page\b", raw))
            result["eof_marker_count"] = raw.count(b"%%EOF")
        except Exception as exc:
            result["metadata_error"] = str(exc)
    if fmt in {"DOCX", "XLSX", "PPTX", "ODT", "ODS", "ODP"}:
        try:
            with zipfile.ZipFile(path) as archive:
                result["zip_member_count"] = len(archive.namelist())
                result["zip_members_preview"] = archive.namelist()[:25]
        except Exception as exc:
            result["zip_probe_error"] = str(exc)
    return result


def raw_docx_paragraph_count(path: Path) -> int | None:
    try:
        with zipfile.ZipFile(path) as archive:
            xml = archive.read("word/document.xml")
        root = ET.fromstring(xml)
        return len(root.findall(".//w:p", WORD_NS))
    except Exception:
        return None


def docx_core_property_lines(document: Any) -> list[str]:
    props = document.core_properties
    fields = [
        "author",
        "category",
        "comments",
        "content_status",
        "created",
        "identifier",
        "keywords",
        "language",
        "last_modified_by",
        "modified",
        "subject",
        "title",
        "version",
    ]
    lines: list[str] = []
    found = False
    for field in fields:
        value = getattr(props, field, None)
        if value not in (None, ""):
            found = True
            lines.append(f"- {field}: {format_value(value)}")
    return lines if found else ["- core_properties: none populated"]


def paragraph_collection_text(paragraphs: Any) -> str:
    values = [paragraph.text.strip() for paragraph in paragraphs if paragraph.text.strip()]
    return repr(values) if values else "none"


def workbook_properties_lines(workbook: Any) -> list[str]:
    props = getattr(workbook, "properties", None)
    if props is None:
        return ["- workbook_properties: unavailable"]
    fields = ["creator", "lastModifiedBy", "title", "subject", "description", "keywords", "category", "created", "modified"]
    lines = ["- workbook_properties:"]
    found = False
    for field in fields:
        value = getattr(props, field, None)
        if value not in (None, ""):
            lines.append(f"  - {field}: {format_value(value)}")
            found = True
    if not found:
        lines.append("  - none populated")
    return lines


def defined_names_lines(workbook: Any) -> list[str]:
    defined_names = getattr(workbook, "defined_names", None)
    if not defined_names:
        return ["- defined_names_count: 0"]
    lines = ["- defined_names:"]
    count = 0
    try:
        iterable = defined_names.items()
    except AttributeError:
        iterable = []
    for name, defined_name in iterable:
        count += 1
        try:
            destinations = [f"{sheet}!{coord}" for sheet, coord in defined_name.destinations]
        except Exception:
            destinations = []
        target = getattr(defined_name, "attr_text", None) or ", ".join(destinations) or "<unknown>"
        lines.append(f"  - name={name!r}; target={target!r}; localSheetId={getattr(defined_name, 'localSheetId', None)!r}")
    if count == 0:
        return ["- defined_names_count: 0"]
    return [f"- defined_names_count: {count}", *lines]


def xlsx_sheet_lines(
    sheet_index: int,
    sheet: Any,
    value_sheet: Any,
    workbook_values: Any,
    sheet_cells: list[Any],
    *,
    truncated: bool = False,
) -> list[str]:
    lines = [
        "",
        f"Worksheet {sheet_index}: {sheet.title}",
        f"- visibility: {sheet.sheet_state}",
        f"- dimensions: {sheet.dimensions}",
        f"- max_row: {sheet.max_row}",
        f"- max_column: {sheet.max_column}",
        f"- freeze_panes: {format_value(sheet.freeze_panes)}",
        f"- merged_ranges_count: {len(list(sheet.merged_cells.ranges))}",
    ]
    for merged in sheet.merged_cells.ranges:
        lines.append(f"  - merged_range: {str(merged)}")
    lines.append(f"- auto_filter: {format_value(sheet.auto_filter.ref)}")
    lines.extend(xlsx_table_lines(sheet))
    lines.append(f"- non_empty_cells_count: {len(sheet_cells)}")
    if truncated:
        lines.append("- extraction_truncated: true")
    if sheet_cells:
        lines.append("- cells:")
    cross_sheet_refs: set[str] = set()
    for cell in sheet_cells:
        lines.append("  - " + xlsx_cell_line(cell, value_sheet[cell.coordinate]))
        if is_formula(cell.value):
            for ref in _CROSS_SHEET_REF_RE.findall(str(cell.value)):
                cross_sheet_refs.add(ref)
    lines.append("- cross_sheet_formula_references: " + (", ".join(sorted(cross_sheet_refs)) if cross_sheet_refs else "none"))
    lines.extend(xlsx_chart_lines(sheet, workbook_values))
    image_count = len(getattr(sheet, "_images", []) or [])
    lines.append(f"- embedded_images_count: {image_count}")
    for index, image in enumerate(getattr(sheet, "_images", []) or [], start=1):
        lines.append(f"  - image {index}: anchor={anchor_text(getattr(image, 'anchor', None))}")
    return lines


def xlsx_table_lines(sheet: Any) -> list[str]:
    tables = list(getattr(sheet, "tables", {}).values())
    lines = [f"- tables_count: {len(tables)}"]
    for table in tables:
        lines.append(
            "  - table: "
            f"name={getattr(table, 'displayName', None)!r}; ref={getattr(table, 'ref', None)!r}; "
            f"headerRowCount={getattr(table, 'headerRowCount', None)!r}; totalsRowCount={getattr(table, 'totalsRowCount', None)!r}"
        )
    return lines


def xlsx_cell_line(cell: Any, cached_cell: Any) -> str:
    formula = str(cell.value) if is_formula(cell.value) else None
    cached_value = cached_cell.value
    pieces = [f"coordinate={cell.coordinate}", f"data_type={cell.data_type!r}", f"value={format_value(cell.value)}"]
    if formula is not None:
        pieces.append(f"formula={formula!r}")
        pieces.append(f"cached_value={format_value(cached_value) if cached_value is not None else '<not available>'}")
    else:
        pieces.append(f"cached_value={format_value(cached_value)}")
    pieces.append(f"number_format={cell.number_format!r}")
    pieces.append(f"hyperlink={cell.hyperlink.target!r}" if cell.hyperlink else "hyperlink=None")
    if cell.comment:
        pieces.append(f"comment_author={cell.comment.author!r}; comment_text={cell.comment.text!r}")
    else:
        pieces.append("comment=None")
    return "; ".join(pieces)


def xlsx_chart_lines(sheet: Any, workbook_values: Any) -> list[str]:
    charts = list(getattr(sheet, "_charts", []) or [])
    lines = [f"- charts_count: {len(charts)}"]
    for index, chart in enumerate(charts, start=1):
        lines.append(f"  - chart {index}:")
        lines.append(f"    - type: {chart.__class__.__name__}")
        lines.append(f"    - title: {format_value(chart_title_text(getattr(chart, 'title', None)))}")
        lines.append(f"    - anchor: {anchor_text(getattr(chart, 'anchor', None))}")
        lines.append(f"    - style: {format_value(getattr(chart, 'style', None))}")
        lines.append(f"    - x_axis_title: {format_value(chart_title_text(getattr(getattr(chart, 'x_axis', None), 'title', None)))}")
        lines.append(f"    - y_axis_title: {format_value(chart_title_text(getattr(getattr(chart, 'y_axis', None), 'title', None)))}")
        series = list(getattr(chart, "ser", []) or [])
        lines.append(f"    - series_count: {len(series)}")
        for series_index, series_item in enumerate(series, start=1):
            cat_ref = data_source_ref(getattr(series_item, "cat", None)) or data_source_ref(getattr(series_item, "xVal", None))
            val_ref = data_source_ref(getattr(series_item, "val", None)) or data_source_ref(getattr(series_item, "yVal", None))
            lines.append(f"    - series {series_index}:")
            lines.append(f"      - title_or_text_ref: {format_value(series_text(getattr(series_item, 'tx', None)))}")
            lines.append(f"      - category_range: {format_value(cat_ref)}")
            lines.append(f"      - value_range: {format_value(val_ref)}")
            if cat_ref:
                lines.append(f"      - category_values: {format_value(values_from_ref(workbook_values, cat_ref))}")
            if val_ref:
                lines.append(f"      - values: {format_value(values_from_ref(workbook_values, val_ref))}")
    return lines


def non_empty_cells(sheet: Any, *, max_cells: int | None = None) -> tuple[list[Any], bool]:
    cells: list[Any] = []
    for row in sheet.iter_rows():
        for cell in row:
            if cell.value not in (None, ""):
                cells.append(cell)
                if max_cells and len(cells) >= max_cells:
                    return cells, True
    return cells, False


def is_formula(value: Any) -> bool:
    return isinstance(value, str) and value.startswith("=")


def chart_title_text(title: Any) -> str | None:
    if title is None:
        return None
    tx = getattr(title, "tx", None)
    if tx is not None:
        rich = getattr(tx, "rich", None)
        if rich is not None:
            fragments: list[str] = []
            for paragraph in getattr(rich, "p", []) or []:
                for run in getattr(paragraph, "r", []) or []:
                    fragments.append(str(getattr(run, "t", "")))
            text = "".join(fragments).strip()
            if text:
                return text
        str_ref = getattr(tx, "strRef", None)
        if str_ref is not None and getattr(str_ref, "f", None):
            return str(str_ref.f)
    return title if isinstance(title, str) else None


def series_text(tx: Any) -> str | None:
    if tx is None:
        return None
    for attr in ("strRef", "numRef"):
        ref = getattr(tx, attr, None)
        if ref is not None and getattr(ref, "f", None):
            return str(ref.f)
    value = getattr(tx, "v", None)
    return str(value) if value is not None else None


def data_source_ref(obj: Any) -> str | None:
    if obj is None:
        return None
    if getattr(obj, "f", None):
        return str(obj.f)
    for attr in ("numRef", "strRef", "multiLvlStrRef"):
        ref = getattr(obj, attr, None)
        if ref is not None and getattr(ref, "f", None):
            return str(ref.f)
    return None


def values_from_ref(workbook_values: Any, ref: str) -> list[Any]:
    try:
        from openpyxl.utils.cell import range_boundaries

        sheet_name, coord = split_sheet_ref(ref)
        ws = workbook_values[sheet_name]
        min_col, min_row, max_col, max_row = range_boundaries(coord.replace("$", ""))
    except Exception:
        return []
    values: list[Any] = []
    for row in ws.iter_rows(min_row=min_row, max_row=max_row, min_col=min_col, max_col=max_col):
        for cell in row:
            values.append(cell.value)
    return values


def split_sheet_ref(ref: str) -> tuple[str, str]:
    if "!" not in ref:
        raise ValueError(ref)
    sheet, coord = ref.rsplit("!", 1)
    sheet = sheet.strip()
    if sheet.startswith("'") and sheet.endswith("'"):
        sheet = sheet[1:-1].replace("''", "'")
    return sheet, coord


def anchor_text(anchor: Any) -> str:
    marker = getattr(anchor, "_from", None)
    if marker is None:
        return format_value(anchor)
    return (
        f"zero_based_col={getattr(marker, 'col', None)}; zero_based_row={getattr(marker, 'row', None)}; "
        f"colOff={getattr(marker, 'colOff', None)}; rowOff={getattr(marker, 'rowOff', None)}"
    )


def slide_notes(slide: Any) -> str:
    try:
        notes_slide = slide.notes_slide
    except Exception:
        return ""
    texts: list[str] = []
    for shape in notes_slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = " ".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
            if text:
                texts.append(text)
    return " | ".join(texts)


def pdf_metadata_text(reader: Any) -> str:
    metadata = reader.metadata or {}
    if not metadata:
        return "none"
    return "; ".join(f"{key}={value!r}" for key, value in metadata.items())


def odf_text_blocks(root: ET.Element, *, max_blocks: int | None = None) -> tuple[list[str], bool]:
    blocks: list[str] = []

    def add(block: str) -> bool:
        if not block:
            return False
        blocks.append(block)
        return bool(max_blocks and len(blocks) >= max_blocks)

    for tag in ("text:p", "text:h"):
        for element in root.findall(f".//{tag}", ODF_NS):
            text = "".join(element.itertext()).strip()
            if add(text):
                return blocks, True
    for table in root.findall(".//table:table", ODF_NS):
        table_name = table.attrib.get(f"{{{ODF_NS['table']}}}name", "table")
        rows: list[str] = []
        for row in table.findall(".//table:table-row", ODF_NS):
            values = ["".join(cell.itertext()).strip() for cell in row.findall("table:table-cell", ODF_NS)]
            values = [value for value in values if value]
            if values:
                rows.append(" | ".join(values))
        if rows:
            if add(f"{table_name}: " + " / ".join(rows)):
                return blocks, True
    return blocks, False


def format_value(value: Any) -> str:
    if value is None:
        return "None"
    if hasattr(value, "isoformat"):
        try:
            return value.isoformat()
        except Exception:
            pass
    return repr(value)


def document_module_params(ctx: ConversionContext, classification: Classification) -> dict[str, Any]:
    cfg = config_for_context(ctx.config)
    return cfg.module_params(classification.rough_category, classification.concrete_format.casefold())


def positive_int(value: Any) -> int | None:
    try:
        number = int(value)
    except (TypeError, ValueError):
        return None
    return number if number > 0 else None


def bool_param(value: Any, default: bool) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "on", "enabled"}
    return bool(value)
